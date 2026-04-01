from __future__ import annotations

from html import escape
from pathlib import Path
import shutil
import subprocess
import zipfile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

TITLE = "The Adversarial Architecture"
PRESENTATION_FOOTER = "Matthew Williams — Senior SWE (MSFT) · 2026-04-01 · Based on Anthropic Labs research by Prithvi Rajasekaran"
CORE_SLIDES = 22
APPENDIX_SLIDES = 5
TOTAL_SLIDES = CORE_SLIDES + APPENDIX_SLIDES
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

NAVY = RGBColor(0x1B, 0x2A, 0x4A)
NAVY_DARK = RGBColor(0x13, 0x20, 0x3A)
TEAL = RGBColor(0x2E, 0xC4, 0xB6)
TEAL_DARK = RGBColor(0x22, 0x9E, 0x92)
CORAL = RGBColor(0xE0, 0x7A, 0x5F)
AMBER = RGBColor(0xF2, 0xCC, 0x8F)
GRAY = RGBColor(0x8B, 0x8B, 0x8B)
GRAY_DARK = RGBColor(0x5B, 0x66, 0x78)
LIGHT_BG = RGBColor(0xF8, 0xF8, 0xF8)
LIGHT_PANEL = RGBColor(0xF0, 0xF0, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SUCCESS_BG = RGBColor(0xE4, 0xF7, 0xF4)
DANGER_BG = RGBColor(0xFA, 0xE9, 0xE4)
WARNING_BG = RGBColor(0xFB, 0xF4, 0xE5)
NEUTRAL_BG = RGBColor(0xEA, 0xEC, 0xF0)
BLUE_GRAY = RGBColor(0x4C, 0x62, 0x7A)
LAVENDER = RGBColor(0xE9, 0xE2, 0xFA)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)

FONT_BODY = "Calibri"
FONT_MONO = "Consolas"
FALLBACK_DARK_SLIDES = {1, 5, 6, 13, 20, 22, 27}

BENCHMARKS = [
    ("small-bugfix-python", "Python", 500, "Fix ± sign bug in arg parser", CORAL),
    ("small-feature-typescript", "TypeScript", 800, "Implement retry with backoff", TEAL),
    ("small-bugfix-go", "Go", 600, "Fix connection pool race condition", BLUE_GRAY),
    ("medium-feature-python", "Python", 1700, "Add tags to FastAPI TODO app", AMBER),
    ("medium-feature-fullstack", "React + Express", 3000, "Real-time notifications", PURPLE),
]

RUN_MATRIX = [
    ("e7c84a5d", "small-bugfix-python", "solo", "1", "PASS", "8.5", "905s"),
    ("bd67944a", "small-bugfix-python", "trio", "3*", "FAIL*", "N/A", "1009s"),
    ("b153e749", "small-bugfix-python", "trio", "2", "PASS", "9.5", "427s"),
    ("efab0ba4", "small-feature-typescript", "solo", "1", "PASS", "8.5", "187s"),
    ("867e4e79", "small-feature-typescript", "trio", "1", "PASS", "8.5", "315s"),
    ("7799434e", "small-bugfix-go", "solo", "1", "FAIL", "6.75", "150s"),
    ("f584e402", "small-bugfix-go", "trio", "3", "FAIL", "7.25", "830s"),
    ("3061e233", "medium-feature-python", "solo", "1", "PASS", "8.5", "297s"),
    ("6649b0bc", "medium-feature-python", "trio", "2", "PASS", "8.0", "1256s"),
    ("410f76ce", "medium-feature-fullstack", "solo", "1", "FAIL", "6.25", "383s"),
    ("7dbac7be", "medium-feature-fullstack", "trio", "1", "PASS", "8.0", "619s"),
]

BACKEND_CRITERIA = [
    ("Product Depth", "HIGH", "7", "Real business logic, edge cases, thoughtful data modeling"),
    ("Functionality", "HIGH", "6", "Core workflow works end-to-end with proper error handling"),
    ("Code Quality", "MEDIUM", "5", "Clean modules, naming, structure, maintainability"),
    ("Test Coverage", "MEDIUM", "5", "Meaningful tests that catch real bugs"),
]

FULLSTACK_CRITERIA = [
    ("Product Depth", "HIGH", "7", "Rich interactive behavior and domain depth"),
    ("Functionality", "HIGH", "6", "UI → API → DB workflow works end-to-end"),
    ("Visual Design", "HIGH", "6", "Coherent layout, spacing, hierarchy, accessibility"),
    ("Code Quality", "MEDIUM", "5", "Component decomposition, API separation, shared patterns"),
]



def set_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color



def apply_text_frame(
    text_frame,
    lines: list[dict[str, object]],
    *,
    margin: float = 0.06,
) -> None:
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.margin_left = Inches(margin)
    text_frame.margin_right = Inches(margin)
    text_frame.margin_top = Inches(margin)
    text_frame.margin_bottom = Inches(margin)
    for index, spec in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        text = str(spec.get("text", ""))
        paragraph.text = text
        paragraph.alignment = spec.get("align", PP_ALIGN.LEFT)
        paragraph.space_after = Pt(spec.get("space_after", 6))
        paragraph.line_spacing = spec.get("line_spacing", 1.1)
        if not paragraph.runs:
            paragraph.add_run()
        for run in paragraph.runs:
            font = run.font
            font.name = spec.get("font", FONT_BODY)
            font.size = Pt(spec.get("size", 18))
            font.bold = spec.get("bold", False)
            font.italic = spec.get("italic", False)
            font.color.rgb = spec.get("color", NAVY)



def add_textbox(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    lines: list[dict[str, object]],
    *,
    margin: float = 0.06,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    apply_text_frame(box.text_frame, lines, margin=margin)
    return box



def add_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor,
    line: RGBColor | None = None,
    rounded: bool = True,
):
    shape_type = (
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        if rounded
        else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    )
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    shape.line.width = Pt(1.5)
    return shape



def add_box_with_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor,
    line: RGBColor | None = None,
    rounded: bool = True,
    lines: list[dict[str, object]],
    margin: float = 0.09,
):
    shape = add_box(slide, x, y, w, h, fill=fill, line=line, rounded=rounded)
    apply_text_frame(shape.text_frame, lines, margin=margin)
    return shape



def add_line(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    *,
    color: RGBColor,
    width: float = 2.0,
):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line



def add_title(slide, number: int, title: str, *, dark: bool = False, subtitle: str | None = None) -> None:
    text_color = WHITE if dark else NAVY
    accent = TEAL if dark else TEAL_DARK
    add_textbox(
        slide,
        0.65,
        0.34,
        10.9,
        0.55,
        [{"text": title, "size": 29, "bold": True, "color": text_color}],
    )
    add_line(slide, 0.65, 0.92, 12.55, 0.92, color=accent, width=2.5)
    if subtitle:
        add_textbox(
            slide,
            0.65,
            0.98,
            10.8,
            0.35,
            [{"text": subtitle, "size": 14, "color": accent}],
        )
    add_textbox(
        slide,
        12.2,
        6.95,
        0.65,
        0.24,
        [{"text": f"{number:02d}", "size": 10, "bold": True, "color": text_color if dark else GRAY_DARK, "align": PP_ALIGN.RIGHT}],
        margin=0,
    )



def add_pill(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    fill: RGBColor,
    color: RGBColor = WHITE,
    size: int = 11,
):
    add_box_with_text(
        slide,
        x,
        y,
        w,
        h,
        fill=fill,
        line=fill,
        lines=[{"text": text, "size": size, "bold": True, "color": color, "align": PP_ALIGN.CENTER}],
        margin=0.02,
    )



def add_quote_strip(slide, quote: str, attribution: str, *, y: float, dark_fill: RGBColor = NAVY_DARK) -> None:
    add_box_with_text(
        slide,
        0.55,
        y,
        12.25,
        0.95,
        fill=dark_fill,
        line=dark_fill,
        lines=[
            {"text": quote, "size": 16, "italic": True, "color": TEAL, "align": PP_ALIGN.CENTER},
            {"text": attribution, "size": 11, "color": WHITE, "align": PP_ALIGN.CENTER, "space_after": 0},
        ],
        margin=0.08,
    )



def add_metric_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    label: str,
    value: str,
    *,
    fill: RGBColor,
    accent: RGBColor,
    text_color: RGBColor = NAVY,
):
    add_box(slide, x, y, w, h, fill=fill, line=accent)
    add_textbox(slide, x + 0.12, y + 0.08, w - 0.24, 0.22, [{"text": label, "size": 11, "bold": True, "color": accent, "align": PP_ALIGN.CENTER}], margin=0)
    add_textbox(slide, x + 0.08, y + 0.32, w - 0.16, h - 0.42, [{"text": value, "size": 20, "bold": True, "color": text_color, "align": PP_ALIGN.CENTER}], margin=0)



def add_benchmark_bar(slide, x: float, y: float, w: float, size: int, color: RGBColor) -> None:
    add_box(slide, x, y + 0.05, w, 0.16, fill=NEUTRAL_BG, line=NEUTRAL_BG, rounded=True)
    fill_w = max(0.4, w * size / 3000)
    add_box(slide, x, y + 0.05, fill_w, 0.16, fill=color, line=color, rounded=True)
    add_textbox(slide, x + w + 0.08, y - 0.02, 0.5, 0.22, [{"text": f"{size}", "size": 11, "bold": True, "color": NAVY}], margin=0)



def add_manual_chart_axes(slide, x: float, y: float, w: float, h: float) -> None:
    add_line(slide, x, y + h, x + w, y + h, color=WHITE, width=1.5)
    add_line(slide, x, y, x, y + h, color=WHITE, width=1.5)
    for score in range(0, 11, 2):
        plot_y = y + h - (h * score / 10)
        add_line(slide, x - 0.05, plot_y, x + w, plot_y, color=GRAY_DARK, width=0.6)
        add_textbox(slide, x - 0.35, plot_y - 0.08, 0.26, 0.16, [{"text": str(score), "size": 10, "color": WHITE, "align": PP_ALIGN.RIGHT}], margin=0)



def add_series_line(slide, x: float, y: float, w: float, h: float, values: list[float], *, color: RGBColor, label: str) -> None:
    x_positions = [x + 0.6, x + w / 2, x + w - 0.4]
    y_positions = [y + h - (h * value / 10) for value in values]
    for idx in range(len(values) - 1):
        add_line(slide, x_positions[idx], y_positions[idx], x_positions[idx + 1], y_positions[idx + 1], color=color, width=2.5)
    for idx, (plot_x, plot_y, value) in enumerate(zip(x_positions, y_positions, values, strict=True)):
        marker = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(plot_x - 0.09), Inches(plot_y - 0.09), Inches(0.18), Inches(0.18))
        marker.fill.solid()
        marker.fill.fore_color.rgb = color
        marker.line.color.rgb = WHITE
        marker.line.width = Pt(1)
        add_textbox(slide, plot_x - 0.18, plot_y - 0.34, 0.36, 0.16, [{"text": f"{value:g}", "size": 9, "bold": True, "color": color, "align": PP_ALIGN.CENTER}], margin=0)
        if idx < 3:
            add_textbox(slide, plot_x - 0.22, y + h + 0.1, 0.44, 0.18, [{"text": f"Iter {idx + 1}", "size": 10, "color": WHITE, "align": PP_ALIGN.CENTER}], margin=0)
    add_pill(slide, x + w - 1.25, y - 0.28, 1.2, 0.22, label, fill=color, color=WHITE, size=10)



def slide_1(slide, number: int) -> None:
    set_background(slide, NAVY)
    ring_left = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(-0.3), Inches(4.2), Inches(3.4), Inches(3.4))
    ring_left.fill.background()
    ring_left.line.color.rgb = TEAL
    ring_left.line.width = Pt(3)
    ring_right = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(10.0), Inches(-0.35), Inches(3.2), Inches(3.2))
    ring_right.fill.background()
    ring_right.line.color.rgb = CORAL
    ring_right.line.width = Pt(3)
    add_line(slide, 2.0, 5.8, 10.9, 1.6, color=AMBER, width=2.2)
    add_textbox(slide, 0.95, 1.55, 8.8, 1.2, [{"text": TITLE, "size": 31, "bold": True, "color": WHITE}])
    add_textbox(slide, 0.98, 2.55, 7.5, 0.45, [{"text": "Why the best AI output comes from agents that disagree", "size": 18, "color": TEAL}])
    add_textbox(slide, 0.98, 6.2, 9.5, 0.28, [{"text": PRESENTATION_FOOTER, "size": 12, "color": WHITE}])
    add_textbox(slide, 12.2, 6.95, 0.65, 0.24, [{"text": f"{number:02d}", "size": 10, "bold": True, "color": WHITE, "align": PP_ALIGN.RIGHT}], margin=0)



def slide_2(slide, number: int) -> None:
    set_background(slide, LIGHT_BG)
    add_title(slide, number, "AI Agents Hit Two Walls")
    add_box(slide, 0.7, 1.35, 7.2, 4.35, fill=WHITE, line=NEUTRAL_BG)
    add_textbox(slide, 1.05, 2.05, 3.3, 0.35, [{"text": "Wall 1 · Context degradation", "size": 18, "bold": True, "color": NAVY}])
    add_textbox(
        slide,
        1.05,
        2.5,
        2.85,
        1.65,
        [{"text": "As the context window fills, coherence drops. Some models start wrapping up work early — even when the job is unfinished.", "size": 16, "color": NAVY, "space_after": 0}],
        margin=0,
    )
    stair_specs = [(3.95, 3.85, 2.15, 0.6), (4.3, 3.15, 2.0, 0.6), (4.65, 2.45, 1.85, 0.6)]
    stair_colors = [CORAL, RGBColor(0xD1, 0x8A, 0x74), RGBColor(0xBE, 0x9A, 0x89)]
    stair_labels = ["Attempt N", "Attempt 5", "Attempt 1"]
    for (x, y, w, h), color, label in zip(stair_specs, stair_colors, stair_labels, strict=True):
        add_box_with_text(slide, x, y, w, h, fill=color, line=color, rounded=False, lines=[{"text": label, "size": 15, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}])
    add_box(slide, 8.15, 1.35, 4.45, 4.35, fill=WHITE, line=NEUTRAL_BG)
    add_textbox(slide, 8.5, 2.05, 3.55, 0.35, [{"text": "Wall 2 · Self-evaluation failure", "size": 18, "bold": True, "color": NAVY}])
    add_textbox(
        slide,
        8.5,
        2.45,
        3.55,
        0.95,
        [{"text": "When agents judge their own output, they reliably talk themselves into approval — even when the quality is obviously mediocre.", "size": 16, "color": NAVY, "space_after": 0}],
        margin=0,
    )
    add_box_with_text(slide, 8.65, 3.2, 3.4, 0.8, fill=LIGHT_PANEL, line=TEAL, lines=[{"text": '“This looks great! Ship it.”', "size": 18, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER}], margin=0.06)
    add_box_with_text(
        slide,
        8.7,
        4.2,
        3.35,
        0.98,
        fill=NAVY_DARK,
        line=NAVY_DARK,
        lines=[
            {"text": "ERROR Traceback (most recent call last):", "size": 13, "color": CORAL, "font": FONT_MONO},
            {"text": "AssertionError: expected 6, got 4", "size": 12, "color": WHITE, "font": FONT_MONO},
        ],
    )
    add_quote_strip(
        slide,
        '“I watched it identify legitimate issues, then talk itself into deciding they weren’t a big deal and approve the work anyway.”',
        "Anthropic Labs",
        y=6.1,
    )



def slide_3(slide, number: int) -> None:
    set_background(slide, LIGHT_BG)
    add_title(slide, number, "What if the Builder and the Critic Were Different Agents?")
    planner = add_box_with_text(slide, 0.9, 1.9, 2.3, 1.4, fill=BLUE_GRAY, line=BLUE_GRAY, lines=[{"text": "Planner\nExpands prompt\ninto spec", "size": 18, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}], margin=0.08)
    generator = add_box_with_text(slide, 4.0, 1.9, 2.5, 1.4, fill=TEAL, line=TEAL, lines=[{"text": "Generator\nBuilds features\nin sprints", "size": 18, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}], margin=0.08)
    evaluator = add_box_with_text(slide, 8.0, 1.9, 2.6, 1.4, fill=CORAL, line=CORAL, lines=[{"text": "Evaluator\nTests and\ngrades", "size": 18, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}], margin=0.08)
    add_line(slide, 3.2, 2.6, 3.95, 2.6, color=GRAY_DARK, width=2.5)
    arrow_right = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(3.55), Inches(2.43), Inches(0.45), Inches(0.34))
    arrow_right.fill.solid(); arrow_right.fill.fore_color.rgb = GRAY_DARK; arrow_right.line.color.rgb = GRAY_DARK
    loop_top = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(6.55), Inches(2.15), Inches(0.55), Inches(0.33))
    loop_top.fill.solid(); loop_top.fill.fore_color.rgb = AMBER; loop_top.line.color.rgb = AMBER
    loop_bottom = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(7.4), Inches(2.72), Inches(0.55), Inches(0.33))
    loop_bottom.fill.solid(); loop_bottom.fill.fore_color.rgb = AMBER; loop_bottom.line.color.rgb = AMBER
    loop_bottom.rotation = 180
    add_textbox(slide, 6.55, 2.45, 1.35, 0.22, [{"text": "feedback loop", "size": 11, "bold": True, "color": AMBER, "align": PP_ALIGN.CENTER}], margin=0)
    add_box_with_text(
        slide,
        0.95,
        4.1,
        11.65,
        1.95,
        fill=WHITE,
        line=NEUTRAL_BG,
        lines=[
            {"text": "• Inspired by GANs: generator builds, evaluator tears it apart, feedback drives improvement", "size": 18, "color": NAVY},
            {"text": "• Separation doesn’t eliminate leniency — but it makes calibration tractable", "size": 18, "color": NAVY},
            {"text": "• Add a Planner so the system builds the right thing before it optimizes the build", "size": 18, "color": NAVY, "space_after": 0},
        ],
    )



def slide_4(slide, number: int) -> None:
    set_background(slide, LIGHT_BG)
    add_title(slide, number, "The Article That Started This")
    headers = [(0.75, 1.55, 2.6, "Approach"), (3.4, 1.55, 2.0, "Duration"), (5.45, 1.55, 1.5, "Cost"), (7.0, 1.55, 5.55, "Result")]
    for x, y, w, text in headers:
        add_box_with_text(slide, x, y, w, 0.42, fill=NAVY, line=NAVY, rounded=False, lines=[{"text": text, "size": 14, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}], margin=0.02)
    rows = [
        ("Solo agent (Opus 4.5)", "20 min", "$9", "Core feature broken, poor layout", DANGER_BG, CORAL),
        ("3-agent harness (Opus 4.5)", "6 hr", "$200", "16-feature app, working core, polished UI", SUCCESS_BG, TEAL_DARK),
        ("Simplified harness (Opus 4.6)", "3 hr 50 min", "$125", "Full DAW, 3 QA rounds", SUCCESS_BG, TEAL),
    ]
    y = 2.0
    for approach, duration, cost, result, fill, accent in rows:
        cells = [(0.75, 2.6, approach), (3.4, 2.0, duration), (5.45, 1.5, cost), (7.0, 5.55, result)]
        for x, w, text in cells:
            add_box_with_text(slide, x, y, w, 0.6, fill=fill, line=WHITE, rounded=False, lines=[{"text": text, "size": 14, "bold": x == 0.75, "color": NAVY, "align": PP_ALIGN.LEFT if x == 0.75 or x == 7.0 else PP_ALIGN.CENTER}])
        add_box(slide, 0.78, y + 0.09, 0.1, 0.42, fill=accent, line=accent, rounded=False)
        y += 0.68
    add_box_with_text(slide, 0.9, 4.55, 11.7, 0.95, fill=AMBER, line=AMBER, lines=[{"text": "Not incrementally better — categorically different.", "size": 24, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER}])
    add_textbox(slide, 0.95, 5.75, 10.5, 0.28, [{"text": "Source: Harness Design for Long-Running Apps — Anthropic Labs (2025)", "size": 12, "color": NAVY}], margin=0)



def slide_5(slide, number: int) -> None:
    set_background(slide, NAVY)
    add_title(slide, number, "The Architecture — Three Agents, One Pipeline", dark=True)
    add_box(slide, 0.8, 1.35, 11.75, 4.8, fill=NAVY_DARK, line=TEAL, rounded=True)
    add_textbox(slide, 4.95, 1.52, 2.5, 0.2, [{"text": "ORCHESTRATOR", "size": 14, "bold": True, "color": TEAL, "align": PP_ALIGN.CENTER}], margin=0)
    add_box_with_text(slide, 1.4, 2.1, 2.2, 1.65, fill=BLUE_GRAY, line=BLUE_GRAY, lines=[{"text": "Planner\n1–4 sentence prompt → ambitious spec\nWHAT, not HOW", "size": 18, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}])
    add_box_with_text(slide, 4.6, 2.1, 2.6, 1.65, fill=TEAL, line=TEAL, lines=[{"text": "Generator\nImplements in sprints\nNegotiates contracts", "size": 18, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}])
    add_box_with_text(slide, 8.15, 2.1, 2.7, 1.65, fill=CORAL, line=CORAL, lines=[{"text": "Evaluator\nTests live app\nGrades against criteria", "size": 18, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}])
    add_line(slide, 3.62, 2.92, 4.55, 2.92, color=WHITE, width=2.5)
    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(4.14), Inches(2.76), Inches(0.38), Inches(0.3))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = WHITE; arrow.line.color.rgb = WHITE
    add_line(slide, 7.22, 2.55, 8.1, 2.55, color=AMBER, width=2.5)
    add_line(slide, 8.1, 3.26, 7.22, 3.26, color=AMBER, width=2.5)
    arrow1 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(7.62), Inches(2.39), Inches(0.36), Inches(0.26))
    arrow1.fill.solid(); arrow1.fill.fore_color.rgb = AMBER; arrow1.line.color.rgb = AMBER
    arrow2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(7.35), Inches(3.1), Inches(0.36), Inches(0.26))
    arrow2.fill.solid(); arrow2.fill.fore_color.rgb = AMBER; arrow2.line.color.rgb = AMBER; arrow2.rotation = 180
    add_pill(slide, 6.98, 2.78, 1.15, 0.24, "adversarial", fill=AMBER, color=NAVY, size=10)
    add_box_with_text(slide, 1.25, 4.45, 10.85, 0.8, fill=GRAY_DARK, line=GRAY_DARK, lines=[{"text": "Telemetry layer · timing · cost · scores · bugs · trends", "size": 18, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}], margin=0.06)
    add_textbox(slide, 1.05, 5.6, 11.2, 0.4, [{"text": "Planner defines the target. Evaluator enforces quality. Telemetry makes the claims measurable.", "size": 15, "color": AMBER, "align": PP_ALIGN.CENTER}], margin=0)



def slide_6(slide, number: int) -> None:
    set_background(slide, NAVY_DARK)
    add_title(slide, number, "How We Prevent Gaming", dark=True)
    add_box_with_text(slide, 0.8, 1.35, 7.95, 4.9, fill=NAVY, line=TEAL, lines=[{"text": "Generator worktree", "size": 18, "bold": True, "color": TEAL}, {"text": "benchmarks/\n├── TASK.md\n├── tests/\n├── src/\n└── _eval/   hidden from generator", "size": 17, "color": WHITE, "font": FONT_MONO}], margin=0.08)
    strike = add_line(slide, 4.2, 3.95, 6.75, 3.95, color=CORAL, width=3)
    add_pill(slide, 5.85, 3.55, 1.3, 0.24, "Goodhart guardrail", fill=CORAL, color=WHITE)
    add_box_with_text(slide, 9.0, 1.75, 3.0, 2.8, fill=SUCCESS_BG, line=TEAL, lines=[{"text": "Evaluator worktree", "size": 18, "bold": True, "color": TEAL_DARK}, {"text": "benchmarks/\n├── TASK.md\n├── tests/\n├── src/\n└── _eval/   visible", "size": 16, "color": NAVY, "font": FONT_MONO}], margin=0.08)
    add_line(slide, 8.65, 1.45, 8.65, 6.0, color=AMBER, width=3)
    add_pill(slide, 8.28, 1.15, 0.74, 0.24, "LOCK", fill=AMBER, color=NAVY)
    add_box_with_text(slide, 9.0, 4.85, 3.0, 1.1, fill=AMBER, line=AMBER, lines=[{"text": "Result\n11 runs · 0 boundary violations", "size": 18, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER}], margin=0.08)
    add_textbox(slide, 1.02, 6.33, 10.5, 0.38, [{"text": "If the generator can see the hidden tests, it optimizes for passing them — not for writing genuinely good code.", "size": 14, "color": WHITE}], margin=0)



def slide_7(slide, number: int) -> None:
    set_background(slide, LIGHT_BG)
    add_title(slide, number, "Building a Critic That Actually Criticizes")
    add_box_with_text(slide, 0.8, 1.25, 12.0, 1.25, fill=LAVENDER, line=PURPLE, lines=[{"text": '“Used an LLM to improve my argument over 4 hours… ask it to argue the opposite. LLM demolishes the entire argument. lol.” — Andrej Karpathy', "size": 17, "italic": True, "color": NAVY, "align": PP_ALIGN.CENTER}], margin=0.08)
    add_box_with_text(slide, 1.0, 3.0, 5.85, 2.4, fill=NAVY_DARK, line=NAVY_DARK, lines=[
        {"text": '{  "functionality": 1,  "verdict": "FAIL" }', "size": 19, "bold": True, "color": CORAL, "font": FONT_MONO, "align": PP_ALIGN.CENTER},
        {"text": '↓ feedback loop', "size": 12, "bold": True, "color": AMBER, "font": FONT_MONO, "align": PP_ALIGN.CENTER},
        {"text": '{  "functionality": 10,  "verdict": "PASS" }', "size": 19, "bold": True, "color": TEAL, "font": FONT_MONO, "align": PP_ALIGN.CENTER},
    ], margin=0.1)
    add_box_with_text(slide, 7.15, 2.8, 5.1, 2.7, fill=WHITE, line=AMBER, lines=[
        {"text": "Anti-people-pleasing rules", "size": 18, "bold": True, "color": NAVY},
        {"text": "• Any test fails → functionality ≤ 4", "size": 16, "color": NAVY},
        {"text": "• No new tests → test coverage ≤ 4", "size": 16, "color": NAVY},
        {"text": "• Stub implementation → product depth ≤ 3", "size": 16, "color": NAVY},
        {"text": "• All scores ≥ 7? Flag as suspicious", "size": 16, "color": NAVY, "space_after": 0},
    ], margin=0.08)
    add_box_with_text(slide, 7.15, 5.7, 5.1, 0.65, fill=WARNING_BG, line=AMBER, lines=[{"text": "Calibration still imperfect: TS benchmark scored func=8 with 50% tests failing.", "size": 13, "color": NAVY, "align": PP_ALIGN.CENTER}], margin=0.04)



def slide_8(slide, number: int) -> None:
    set_background(slide, LIGHT_BG)
    add_title(slide, number, "Agreeing on “Done” Before Writing Code")
    add_box_with_text(slide, 0.85, 1.55, 4.2, 3.7, fill=WHITE, line=TEAL_DARK, lines=[
        {"text": "Generator proposal", "size": 20, "bold": True, "color": TEAL_DARK, "align": PP_ALIGN.CENTER},
        {"text": "• Implement the next feature slice", "size": 16, "color": NAVY},
        {"text": "• Define concrete acceptance checks", "size": 16, "color": NAVY},
        {"text": "• Commit to what will be demo-ready this sprint", "size": 16, "color": NAVY},
    ], margin=0.08)
    add_box_with_text(slide, 8.3, 1.55, 4.2, 3.7, fill=WHITE, line=CORAL, lines=[
        {"text": "Evaluator review", "size": 20, "bold": True, "color": CORAL, "align": PP_ALIGN.CENTER},
        {"text": "• Tighten or remove vague criteria", "size": 16, "color": NAVY},
        {"text": "• Add missing failure cases", "size": 16, "color": NAVY},
        {"text": "• Cap negotiation at two rounds", "size": 16, "color": NAVY},
    ], margin=0.08)
    add_line(slide, 5.15, 3.05, 6.15, 3.05, color=GRAY_DARK, width=2.2)
    add_line(slide, 7.1, 3.05, 8.2, 3.05, color=GRAY_DARK, width=2.2)
    add_box_with_text(slide, 5.55, 2.45, 1.95, 1.2, fill=AMBER, line=AMBER, lines=[{"text": "Agreed\ncontract", "size": 20, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER}], margin=0.08)
    add_box_with_text(slide, 2.4, 5.75, 8.65, 0.7, fill=LIGHT_PANEL, line=LIGHT_PANEL, lines=[{"text": "Why it matters: the Planner stays high-level, but each sprint still lands on a testable definition of done.", "size": 16, "color": NAVY, "align": PP_ALIGN.CENTER}], margin=0.06)
    add_pill(slide, 5.55, 1.65, 1.95, 0.24, "Article example: 27 acceptance criteria in one sprint", fill=NAVY, color=WHITE, size=10)



def slide_9(slide, number: int) -> None:
    set_background(slide, LIGHT_BG)
    add_title(slide, number, "Files on Disk, Not Chat History")
    add_box_with_text(slide, 0.8, 1.45, 5.35, 4.95, fill=WHITE, line=NEUTRAL_BG, lines=[
        {"text": "Copilot skill artifacts", "size": 19, "bold": True, "color": NAVY},
        {"text": "harnessa-spec.md\nharnessa-gen-report.md\nharnessa-eval.md", "size": 17, "color": TEAL_DARK, "font": FONT_MONO},
        {"text": "", "size": 4},
        {"text": "Benchmark harness", "size": 19, "bold": True, "color": NAVY},
        {"text": "planner/spec.md\ncontracts/\nevaluations/\ntelemetry/", "size": 17, "color": CORAL, "font": FONT_MONO},
    ], margin=0.08)
    add_box_with_text(slide, 6.5, 1.65, 5.75, 3.8, fill=WHITE, line=TEAL, lines=[
        {"text": "Why files win", "size": 20, "bold": True, "color": NAVY},
        {"text": "• Full audit trail", "size": 17, "color": NAVY},
        {"text": "• Human-readable handoffs", "size": 17, "color": NAVY},
        {"text": "• Resumability after crashes", "size": 17, "color": NAVY},
        {"text": "• Clean phase boundaries between agents", "size": 17, "color": NAVY, "space_after": 0},
    ], margin=0.08)
    add_line(slide, 2.35, 3.0, 4.55, 3.0, color=AMBER, width=2.6)
    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(4.2), Inches(2.83), Inches(0.3), Inches(0.28))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = AMBER; arrow.line.color.rgb = AMBER
    add_line(slide, 4.95, 3.0, 5.8, 3.0, color=AMBER, width=2.6)
    add_box_with_text(slide, 6.5, 5.75, 5.75, 0.7, fill=GRAY_DARK, line=GRAY_DARK, lines=[{"text": "Every write is atomic — crash-safe by design.", "size": 16, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}], margin=0.05)



def slide_10(slide, number: int) -> None:
    set_background(slide, LIGHT_BG)
    add_title(slide, number, "5 Benchmarks, 2 Modes, Real Code")
    add_box(slide, 0.7, 1.3, 7.55, 4.7, fill=WHITE, line=NEUTRAL_BG)
    add_textbox(slide, 0.95, 1.55, 7.0, 0.25, [{"text": "Benchmark", "size": 13, "bold": True, "color": NAVY}], margin=0)
    add_textbox(slide, 4.25, 1.55, 1.2, 0.25, [{"text": "Lang", "size": 13, "bold": True, "color": NAVY}], margin=0)
    add_textbox(slide, 5.45, 1.55, 2.4, 0.25, [{"text": "Challenge", "size": 13, "bold": True, "color": NAVY}], margin=0)
    y = 1.9
    for name, lang, size, challenge, color in BENCHMARKS:
        add_line(slide, 0.9, y - 0.05, 7.95, y - 0.05, color=NEUTRAL_BG, width=1)
        add_textbox(slide, 0.95, y + 0.02, 3.2, 0.32, [{"text": name, "size": 13, "bold": True, "color": NAVY}], margin=0)
        add_pill(slide, 4.2, y + 0.02, 1.0 if len(lang) < 8 else 1.45, 0.24, lang, fill=color if color != AMBER else TEAL_DARK, color=WHITE, size=10)
        add_textbox(slide, 5.45, y + 0.02, 2.2, 0.32, [{"text": challenge, "size": 12, "color": NAVY}], margin=0)
        y += 0.82
    add_box_with_text(slide, 8.55, 1.35, 3.95, 4.1, fill=WHITE, line=NEUTRAL_BG, lines=[{"text": "Repo size (LOC)", "size": 18, "bold": True, "color": NAVY}], margin=0.08)
    y = 2.0
    for name, _lang, size, _challenge, color in BENCHMARKS:
        add_textbox(slide, 8.8, y - 0.02, 2.5, 0.18, [{"text": name.replace("small-", "").replace("medium-", ""), "size": 10, "color": NAVY}], margin=0)
        add_benchmark_bar(slide, 8.8, y + 0.18, 2.65, size, color)
        y += 0.72
    add_box_with_text(slide, 0.95, 6.2, 11.6, 0.65, fill=GRAY_DARK, line=GRAY_DARK, lines=[{"text": "Same model · same prompt · same tools · hidden _eval tests + fixtures · only the architecture changes", "size": 15, "color": WHITE, "align": PP_ALIGN.CENTER}], margin=0.05)



def slide_11(slide, number: int) -> None:
    set_background(slide, LIGHT_BG)
    add_title(slide, number, "Solo FAIL → Trio PASS")
    add_box_with_text(slide, 0.8, 1.1, 11.95, 0.42, fill=WARNING_BG, line=AMBER, lines=[{"text": "N=11 total runs — small sample, but the fullstack benchmark produced a categorical result, not a marginal score bump.", "size": 12, "color": NAVY, "align": PP_ALIGN.CENTER}], margin=0.03)
    add_box_with_text(slide, 0.95, 1.8, 5.45, 3.85, fill=DANGER_BG, line=CORAL, lines=[
        {"text": "Solo", "size": 24, "bold": True, "color": CORAL, "align": PP_ALIGN.CENTER},
        {"text": "FAIL", "size": 26, "bold": True, "color": CORAL, "align": PP_ALIGN.CENTER},
        {"text": "Functionality 4 / 10", "size": 18, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER},
        {"text": "Average score 6.25", "size": 17, "color": NAVY, "align": PP_ALIGN.CENTER},
        {"text": "Duration 383s", "size": 17, "color": NAVY, "align": PP_ALIGN.CENTER},
        {"text": "Broken WebSocket notifications", "size": 16, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER},
    ], margin=0.08)
    add_box_with_text(slide, 6.95, 1.8, 5.45, 3.85, fill=SUCCESS_BG, line=TEAL, lines=[
        {"text": "Trio", "size": 24, "bold": True, "color": TEAL_DARK, "align": PP_ALIGN.CENTER},
        {"text": "PASS", "size": 26, "bold": True, "color": TEAL_DARK, "align": PP_ALIGN.CENTER},
        {"text": "Functionality 8 / 10", "size": 18, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER},
        {"text": "Average score 8.0", "size": 17, "color": NAVY, "align": PP_ALIGN.CENTER},
        {"text": "Duration 619s (1.6×)", "size": 17, "color": NAVY, "align": PP_ALIGN.CENTER},
        {"text": "Working WebSocket notifications", "size": 16, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER},
    ], margin=0.08)
    add_textbox(slide, 5.95, 3.1, 1.4, 0.38, [{"text": "+4 functionality", "size": 16, "bold": True, "color": AMBER, "align": PP_ALIGN.CENTER}], margin=0)
    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(5.95), Inches(3.45), Inches(1.15), Inches(0.52))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = AMBER; arrow.line.color.rgb = AMBER
    add_pill(slide, 7.25, 1.45, 1.8, 0.24, "84s planner spec", fill=AMBER, color=NAVY, size=10)
    add_box_with_text(slide, 1.1, 6.0, 11.3, 0.7, fill=NAVY, line=NAVY, lines=[{"text": "This is the categorical difference the article predicted: same model, same task, same tools — one architecture ships broken code, the other ships a working core feature.", "size": 15, "color": WHITE, "align": PP_ALIGN.CENTER}], margin=0.05)



def slide_12(slide, number: int) -> None:
    set_background(slide, LIGHT_BG)
    add_title(slide, number, "Trio Won the Tests That Mattered")
    rows = [
        ("Python bugfix", "PASS 8.5", "PASS 9.5", "Trio", "Evaluator caught issue, gen fixed it", TEAL),
        ("TS feature", "PASS 8.5", "PASS 8.5", "Tie", "Simple task — trio overhead wasted", GRAY),
        ("Go race", "FAIL 6.75", "FAIL 7.25", "Tie", "Too hard for both", GRAY),
        ("Python tags", "PASS 8.5", "PASS 8.0", "Solo*", "3.25 → 8.0 after feedback; solo likely lenient", CORAL),
        ("Fullstack", "FAIL 6.25", "PASS 8.0", "Trio", "Categorical difference", TEAL),
    ]
    y_positions = [1.4, 2.15, 2.9, 3.65, 4.5]
    heights = [0.6, 0.6, 0.6, 0.6, 1.0]
    for (label, solo, trio, winner, insight, accent), y, height in zip(rows, y_positions, heights, strict=True):
        fill = SUCCESS_BG if accent == TEAL else DANGER_BG if accent == CORAL else WHITE
        line = accent if winner != "Tie" else GRAY
        add_box(slide, 0.8, y, 11.8, height, fill=fill if label == "Fullstack" else WHITE, line=line)
        add_textbox(slide, 1.02, y + 0.11, 2.2, 0.3, [{"text": label, "size": 16 if label == "Fullstack" else 14, "bold": True, "color": NAVY}], margin=0)
        add_pill(slide, 3.45, y + 0.13, 1.35, 0.24, solo, fill=CORAL if "FAIL" in solo else NAVY, color=WHITE, size=10)
        add_pill(slide, 5.1, y + 0.13, 1.35, 0.24, trio, fill=TEAL if "PASS" in trio else NAVY, color=WHITE, size=10)
        add_pill(slide, 6.85, y + 0.13, 1.0, 0.24, winner, fill=accent if winner != "Tie" else GRAY, color=WHITE, size=10)
        add_textbox(slide, 8.05, y + 0.09, 3.95, height - 0.1, [{"text": insight, "size": 12 if label != "Fullstack" else 14, "bold": label == "Fullstack", "color": NAVY}], margin=0)
    add_metric_card(slide, 0.95, 5.85, 2.5, 0.78, "Mean functionality", "4.8 → 7.6", fill=WHITE, accent=AMBER)
    add_metric_card(slide, 3.75, 5.85, 2.2, 0.78, "Verdicts", "3/5 → 4/5", fill=WHITE, accent=TEAL)
    add_metric_card(slide, 6.25, 5.85, 2.0, 0.78, "Wall-clock", "~1.8×", fill=WHITE, accent=GRAY)
    add_box_with_text(slide, 8.55, 5.8, 4.0, 0.88, fill=WARNING_BG, line=AMBER, lines=[{"text": "Python tags caveat: solo’s higher numeric score likely reflects self-evaluation leniency, not stronger engineering.", "size": 12, "color": NAVY}], margin=0.05)



def slide_13(slide, number: int) -> None:
    set_background(slide, NAVY)
    add_title(slide, number, "The Feedback Loop Works", dark=True)
    chart_x, chart_y, chart_w, chart_h = 1.0, 1.7, 7.2, 3.9
    add_manual_chart_axes(slide, chart_x, chart_y, chart_w, chart_h)
    threshold_y = chart_y + chart_h - (chart_h * 7 / 10)
    threshold = add_line(slide, chart_x, threshold_y, chart_x + chart_w, threshold_y, color=AMBER, width=1.8)
    threshold.line.dash_style = 4
    add_textbox(slide, 7.45, threshold_y - 0.13, 0.7, 0.2, [{"text": "PASS", "size": 10, "bold": True, "color": AMBER}], margin=0)
    add_series_line(slide, chart_x, chart_y, chart_w, chart_h, [5.0, 9.5, 9.5], color=WHITE, label="Python bugfix")
    add_series_line(slide, chart_x, chart_y, chart_w, chart_h, [2.75, 6.5, 7.25], color=CORAL, label="Go race")
    add_series_line(slide, chart_x, chart_y, chart_w, chart_h, [3.25, 8.0, 8.0], color=TEAL, label="Python tags")
    add_box_with_text(slide, 8.75, 1.75, 3.7, 3.75, fill=NAVY_DARK, line=TEAL, lines=[
        {"text": "What the loop reveals", "size": 20, "bold": True, "color": WHITE},
        {"text": "• Iteration 1 is intentionally harsh", "size": 16, "color": WHITE},
        {"text": "• Scores jump after specific feedback", "size": 16, "color": WHITE},
        {"text": "• Go race plateaus: 3 iterations still not enough", "size": 16, "color": WHITE},
        {"text": "• The article’s core claim holds: improve, then plateau", "size": 16, "color": AMBER, "bold": True, "space_after": 0},
    ], margin=0.08)



def slide_14(slide, number: int) -> None:
    set_background(slide, LIGHT_BG)
    add_title(slide, number, "5 Confirmed, 2 Partial, 2 Inconclusive")
    add_box_with_text(slide, 0.8, 1.5, 4.0, 4.8, fill=SUCCESS_BG, line=TEAL, lines=[
        {"text": "Confirmed", "size": 22, "bold": True, "color": TEAL_DARK, "align": PP_ALIGN.CENTER},
        {"text": "Separating generator and evaluator is a strong lever", "size": 15, "color": NAVY},
        {"text": "Scores improve over iterations", "size": 15, "color": NAVY},
        {"text": "Evaluator worth cost only beyond solo capability", "size": 15, "color": NAVY},
        {"text": "Solo agents have self-evaluation failure", "size": 15, "color": NAVY},
        {"text": "Claude is a poor QA agent out of the box", "size": 15, "color": NAVY, "space_after": 0},
    ], margin=0.08)
    add_box_with_text(slide, 4.95, 1.5, 3.7, 4.8, fill=WARNING_BG, line=AMBER, lines=[
        {"text": "Partial", "size": 22, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER},
        {"text": "Harness output is categorically different — but only on complex tasks", "size": 16, "color": NAVY},
        {"text": "Planner adds more structure than scope expansion", "size": 16, "color": NAVY, "space_after": 0},
    ], margin=0.08)
    add_box_with_text(slide, 8.85, 1.5, 3.75, 4.8, fill=NEUTRAL_BG, line=GRAY, lines=[
        {"text": "Inconclusive", "size": 22, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER},
        {"text": "Criteria wording steers output", "size": 16, "color": NAVY},
        {"text": "Harness assumptions go stale with new models", "size": 16, "color": NAVY, "space_after": 0},
    ], margin=0.08)
    add_box_with_text(slide, 1.2, 6.25, 11.0, 0.45, fill=NAVY, line=NAVY, lines=[{"text": "Bottom line: the article mostly holds up, but the strongest signal appears exactly where the task is hard enough to break a solo agent.", "size": 14, "color": WHITE, "align": PP_ALIGN.CENTER}], margin=0.03)



def slide_15(slide, number: int) -> None:
    set_background(slide, LIGHT_BG)
    add_title(slide, number, "We’re Not Alone — This Is an Industry Pattern")
    add_line(slide, 1.0, 3.45, 12.1, 3.45, color=GRAY_DARK, width=2.2)
    entries = [
        (1.05, 2.1, "Jun 2025", "Anthropic research system", "Lead agent + subagents\n90.2% better than single-agent", TEAL),
        (3.35, 4.0, "Feb 2026", "Anthropic C compiler", "16 parallel agents\n“Harness matters more than agents”", CORAL),
        (5.75, 2.1, "Mar 2026", "Harness Design article", "Generator ↔ evaluator loop\nplanner expands scope", NAVY),
        (8.2, 4.0, "2026", "Claude Code coordinator", "“Do not rubber-stamp weak work”\nanti-sycophancy ships in product", PURPLE),
        (10.55, 2.1, "2026", "OpenAI Symphony", "Issue tracker → isolated workspaces\nworks best with harness engineering", TEAL_DARK),
    ]
    for x, y, date, name, summary, color in entries:
        add_box_with_text(slide, x, y, 1.95, 1.05, fill=WHITE, line=color, lines=[{"text": date, "size": 11, "bold": True, "color": color, "align": PP_ALIGN.CENTER}, {"text": name, "size": 13, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER}, {"text": summary, "size": 11, "color": NAVY, "align": PP_ALIGN.CENTER}], margin=0.05)
        add_line(slide, x + 0.98, 3.15 if y < 3 else 4.0, x + 0.98, 3.45, color=color, width=1.5)
    add_box_with_text(slide, 0.95, 5.45, 11.65, 0.8, fill=NAVY, line=NAVY, lines=[{"text": "Convergence thesis: isolation, structured handoffs, specialized roles, and evaluation as architecture are showing up independently across the major labs.", "size": 16, "color": WHITE, "align": PP_ALIGN.CENTER}], margin=0.06)



def slide_16(slide, number: int) -> None:
    set_background(slide, LIGHT_BG)
    add_title(slide, number, "The Community Is Already Here")
    center = add_box_with_text(slide, 5.1, 2.5, 3.1, 1.2, fill=NAVY, line=NAVY, lines=[{"text": "Adversarial pattern\nDisagreement is signal", "size": 21, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}], margin=0.08)
    satellites = [
        (1.0, 1.65, 2.6, 1.0, "GStack", " /review → /codex\nindependent second opinion", TEAL),
        (1.1, 4.4, 2.6, 0.95, "ClawCompany", "productized harness workflow", CORAL),
        (9.55, 1.65, 2.4, 0.95, "OpenHands", "agent SDK + orchestration", BLUE_GRAY),
        (9.45, 4.4, 2.55, 0.95, "Cross-model review", "Claude + Codex disagreements\nflaged, reconciled, escalated", PURPLE),
    ]
    for x, y, w, h, title, body, color in satellites:
        add_box_with_text(slide, x, y, w, h, fill=WHITE, line=color, lines=[{"text": title, "size": 17, "bold": True, "color": color, "align": PP_ALIGN.CENTER}, {"text": body, "size": 12, "color": NAVY, "align": PP_ALIGN.CENTER}], margin=0.05)
        add_line(slide, x + w / 2, y + h / 2, 6.65, 3.1, color=color, width=1.7)
    add_box_with_text(slide, 3.85, 5.75, 5.65, 0.75, fill=WARNING_BG, line=AMBER, lines=[{"text": "What this means: the pattern is production-ready. GStack already uses the “outside voice” workflow daily in YC companies.", "size": 14, "color": NAVY, "align": PP_ALIGN.CENTER}], margin=0.05)



def slide_17(slide, number: int) -> None:
    set_background(slide, LIGHT_BG)
    add_title(slide, number, "What the Evaluator Catches (That Solo Would Ship)")
    add_box_with_text(slide, 0.9, 1.55, 4.6, 4.3, fill=DANGER_BG, line=CORAL, lines=[
        {"text": "Iteration 1 — rejected", "size": 22, "bold": True, "color": CORAL, "align": PP_ALIGN.CENTER},
        {"text": "1,206-line index.ts monolith", "size": 18, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER},
        {"text": "CDN React + Babel string template", "size": 16, "color": NAVY, "align": PP_ALIGN.CENTER},
        {"text": "0 tests · 0 real components", "size": 16, "color": NAVY, "align": PP_ALIGN.CENTER},
        {"text": "A solo agent would have shipped this.", "size": 18, "italic": True, "color": NAVY, "align": PP_ALIGN.CENTER},
    ], margin=0.08)
    add_box_with_text(slide, 7.05, 1.55, 4.6, 4.3, fill=SUCCESS_BG, line=TEAL, lines=[
        {"text": "Iteration 2 — passed", "size": 22, "bold": True, "color": TEAL_DARK, "align": PP_ALIGN.CENTER},
        {"text": "32-file rebuild", "size": 18, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER},
        {"text": "React + Vite + Tailwind + React Router", "size": 16, "color": NAVY, "align": PP_ALIGN.CENTER},
        {"text": "6 components · 4 pages · 7 tests", "size": 16, "color": NAVY, "align": PP_ALIGN.CENTER},
        {"text": "Spec survived; architecture got rebuilt.", "size": 18, "italic": True, "color": NAVY, "align": PP_ALIGN.CENTER},
    ], margin=0.08)
    add_box_with_text(slide, 5.45, 2.55, 1.15, 2.15, fill=AMBER, line=AMBER, lines=[{"text": "+4,335\n−2,973\n36 file\nchanges", "size": 18, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER}], margin=0.08)
    add_box_with_text(slide, 1.1, 6.1, 10.9, 0.5, fill=NAVY, line=NAVY, lines=[{"text": "The evaluator turned a technically-working monolith into a production-shape rebuild — exactly the kind of upgrade solo self-evaluation misses.", "size": 14, "color": WHITE, "align": PP_ALIGN.CENTER}], margin=0.04)



def slide_18(slide, number: int) -> None:
    set_background(slide, LIGHT_BG)
    add_title(slide, number, "It’s Not Always Worth It")
    add_box_with_text(slide, 0.9, 1.6, 5.4, 3.95, fill=SUCCESS_BG, line=TEAL, lines=[
        {"text": "Use Trio", "size": 24, "bold": True, "color": TEAL_DARK, "align": PP_ALIGN.CENTER},
        {"text": "• Features touching 3+ files", "size": 17, "color": NAVY},
        {"text": "• Unfamiliar bugs with unclear root cause", "size": 17, "color": NAVY},
        {"text": "• Architecture refactors", "size": 17, "color": NAVY},
        {"text": "• Full-stack features", "size": 17, "color": NAVY},
        {"text": "• Anything where “does it actually work?” is uncertain", "size": 17, "color": NAVY, "space_after": 0},
    ], margin=0.08)
    add_box_with_text(slide, 7.0, 1.6, 5.4, 3.95, fill=DANGER_BG, line=CORAL, lines=[
        {"text": "Skip It", "size": 24, "bold": True, "color": CORAL, "align": PP_ALIGN.CENTER},
        {"text": "• Single-line fixes", "size": 17, "color": NAVY},
        {"text": "• Formatting or linting", "size": 17, "color": NAVY},
        {"text": "• Dependency bumps", "size": 17, "color": NAVY},
        {"text": "• Tasks where you already know the fix", "size": 17, "color": NAVY},
        {"text": "• Simple, well-understood changes", "size": 17, "color": NAVY, "space_after": 0},
    ], margin=0.08)
    add_box_with_text(slide, 1.5, 6.0, 10.3, 0.7, fill=AMBER, line=AMBER, lines=[{"text": "Rule of thumb: if a solo agent gets it right on the first try, trio is just overhead. The sweet spot is medium-complexity work at the edge of model capability.", "size": 16, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER}], margin=0.05)



def slide_19(slide, number: int) -> None:
    set_background(slide, LIGHT_BG)
    add_title(slide, number, "What We Found That the Article Didn’t Cover")
    cards = [
        (1.0, "1", "Planner value = structure, not scope", "The spec gave the generator a roadmap. On medium tasks that structure enabled correct first attempts; on small tasks it was overhead."),
        (4.35, "2", "Evaluator leniency persists", "Even with hard rules, we still saw func=8 with 50% tests failing. Guardrails help, but skepticism remains an active calibration problem."),
        (7.7, "3", "Telemetry keeps the claims honest", "Every run emits structured JSON. No “it feels better” hand-waving — if you can’t measure it, you can’t trust it."),
    ]
    for x, n, title, body in cards:
        add_box_with_text(slide, x, 1.85, 3.15, 4.35, fill=WHITE, line=NEUTRAL_BG, lines=[{"text": n, "size": 28, "bold": True, "color": AMBER, "align": PP_ALIGN.CENTER}, {"text": title, "size": 19, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER}, {"text": body, "size": 14, "color": NAVY, "align": PP_ALIGN.CENTER}], margin=0.08)



def slide_20(slide, number: int) -> None:
    set_background(slide, NAVY)
    add_title(slide, number, "One Command, Three Agents", dark=True)
    add_box_with_text(slide, 0.95, 1.45, 11.3, 1.1, fill=NAVY_DARK, line=TEAL, lines=[{"text": "copilot -p '/harnessa Fix the authentication bug' --allow-all", "size": 24, "bold": True, "color": WHITE, "font": FONT_MONO, "align": PP_ALIGN.CENTER}], margin=0.08)
    steps = [
        (1.0, TEAL, "Planner", "Reads the code\nWrites harnessa-spec.md"),
        (4.45, AMBER, "Generator", "Reads only the spec\nWrites harnessa-gen-report.md"),
        (7.95, CORAL, "Evaluator", "Reads report + diff + hidden tests\nWrites harnessa-eval.md"),
    ]
    for x, color, title, body in steps:
        add_box_with_text(slide, x, 3.05, 3.0, 1.85, fill=color, line=color, lines=[{"text": title, "size": 22, "bold": True, "color": NAVY if color == AMBER else WHITE, "align": PP_ALIGN.CENTER}, {"text": body, "size": 15, "color": NAVY if color == AMBER else WHITE, "align": PP_ALIGN.CENTER}], margin=0.08)
    add_line(slide, 4.0, 3.96, 4.42, 3.96, color=WHITE, width=2.2)
    add_line(slide, 7.43, 3.96, 7.92, 3.96, color=WHITE, width=2.2)
    for x in (4.12, 7.55):
        arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(x), Inches(3.8), Inches(0.28), Inches(0.28))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = WHITE; arrow.line.color.rgb = WHITE
    add_box_with_text(slide, 1.15, 5.55, 6.0, 0.55, fill=GRAY_DARK, line=GRAY_DARK, lines=[{"text": "Context resets between phases. Each agent gets a clean slate.", "size": 15, "color": WHITE, "align": PP_ALIGN.CENTER}], margin=0.04)
    add_box_with_text(slide, 7.55, 5.55, 4.7, 0.55, fill=WARNING_BG, line=AMBER, lines=[{"text": "Fallback: use a pre-recorded terminal walkthrough. Never depend on a live API call in a 15-minute slot.", "size": 13, "color": NAVY, "align": PP_ALIGN.CENTER}], margin=0.04)
    add_textbox(slide, 1.05, 6.45, 6.0, 0.24, [{"text": "Open source · MIT · benchmarks + telemetry included · replay mode", "size": 12, "color": TEAL}], margin=0)



def slide_21(slide, number: int) -> None:
    set_background(slide, LIGHT_BG)
    add_title(slide, number, "When Should You Use This?")
    add_box_with_text(slide, 0.9, 1.35, 3.25, 1.85, fill=WHITE, line=NEUTRAL_BG, lines=[{"text": "Decision rule", "size": 20, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER}, {"text": "If a solo agent gets it right first try → don’t add overhead.\nIf the task sits at the edge of capability → trio is insurance.", "size": 15, "color": NAVY, "align": PP_ALIGN.CENTER}], margin=0.08)
    add_box_with_text(slide, 4.55, 1.35, 3.8, 1.85, fill=WARNING_BG, line=AMBER, lines=[{"text": "Model tiering", "size": 20, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER}, {"text": "Opus for plan + evaluation\nSonnet for building volume", "size": 16, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER}, {"text": "Industry practice — not tested in our experiments", "size": 11, "italic": True, "color": NAVY, "align": PP_ALIGN.CENTER}], margin=0.08)
    add_box_with_text(slide, 8.75, 1.35, 3.55, 1.85, fill=NEUTRAL_BG, line=GRAY, lines=[{"text": "Round-robin idea", "size": 20, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER}, {"text": "Claude plans\nGPT builds\nGemini evaluates", "size": 16, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER}, {"text": "Conjecture — immediate next experiment, not validated here", "size": 11, "italic": True, "color": NAVY, "align": PP_ALIGN.CENTER}], margin=0.08)
    add_box_with_text(slide, 1.15, 3.75, 4.2, 2.0, fill=SUCCESS_BG, line=TEAL, lines=[{"text": "Solo", "size": 24, "bold": True, "color": TEAL_DARK, "align": PP_ALIGN.CENTER}, {"text": "Known fix\nsmall change\nlow downside if wrong", "size": 16, "color": NAVY, "align": PP_ALIGN.CENTER}], margin=0.08)
    add_box_with_text(slide, 8.0, 3.75, 4.2, 2.0, fill=SUCCESS_BG, line=CORAL, lines=[{"text": "Trio", "size": 24, "bold": True, "color": CORAL, "align": PP_ALIGN.CENTER}, {"text": "Ambiguous bug\nmulti-file feature\nhigh cost if broken", "size": 16, "color": NAVY, "align": PP_ALIGN.CENTER}], margin=0.08)
    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(5.75), Inches(4.35), Inches(1.7), Inches(0.75))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = AMBER; arrow.line.color.rgb = AMBER
    add_textbox(slide, 5.72, 4.58, 1.75, 0.18, [{"text": "complex?", "size": 13, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER}], margin=0)
    add_box_with_text(slide, 1.0, 6.15, 11.5, 0.45, fill=NAVY, line=NAVY, lines=[{"text": "Bigger question: as models improve, which parts of the harness are still load-bearing? The architecture that adapts wins.", "size": 15, "color": WHITE, "align": PP_ALIGN.CENTER}], margin=0.03)



def slide_22(slide, number: int) -> None:
    set_background(slide, NAVY)
    add_title(slide, number, "The Quality Ceiling is Real. Architecture Breaks Through It.", dark=True)
    add_box_with_text(slide, 1.2, 1.65, 10.9, 1.55, fill=NAVY_DARK, line=TEAL, lines=[
        {"text": "When your solo agent ships broken code, the problem is not just the model — it’s the architecture around the model.", "size": 20, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER},
        {"text": "Separating building from evaluation is a proven, measurable lever.", "size": 18, "color": TEAL, "align": PP_ALIGN.CENTER},
    ], margin=0.08)
    add_box_with_text(slide, 1.35, 3.75, 10.6, 1.65, fill=NAVY, line=NAVY, lines=[{"text": '“The space of interesting harness combinations doesn’t shrink as models improve. It moves.”', "size": 24, "italic": True, "color": AMBER, "align": PP_ALIGN.CENTER}, {"text": "Anthropic Labs", "size": 13, "color": WHITE, "align": PP_ALIGN.CENTER}], margin=0.08)
    add_textbox(slide, 1.15, 6.3, 10.5, 0.28, [{"text": "github.com/ridermw/harnessa · MIT License · Benchmarks + telemetry included", "size": 12, "color": WHITE, "align": PP_ALIGN.CENTER}], margin=0)



def slide_a1(slide, number: int) -> None:
    set_background(slide, LIGHT_BG)
    add_title(slide, number, "Appendix — Full Experiment Matrix")
    headers = ["Run ID", "Benchmark", "Mode", "Iter", "Verdict", "Avg", "Duration"]
    x_positions = [0.7, 2.0, 6.1, 7.2, 8.0, 9.25, 10.4]
    widths = [1.15, 4.0, 0.85, 0.65, 1.0, 0.85, 1.3]
    for x, w, header in zip(x_positions, widths, headers, strict=True):
        add_box_with_text(slide, x, 1.3, w, 0.38, fill=NAVY, line=NAVY, rounded=False, lines=[{"text": header, "size": 11, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}], margin=0.02)
    y = 1.72
    row_h = 0.38
    for idx, row in enumerate(RUN_MATRIX):
        fill = WHITE if idx % 2 == 0 else LIGHT_PANEL
        for x, w, value in zip(x_positions, widths, row, strict=True):
            color = NAVY
            if value == "PASS":
                color = TEAL_DARK
            elif "FAIL" in value:
                color = CORAL
            add_box_with_text(slide, x, y, w, row_h, fill=fill, line=WHITE, rounded=False, lines=[{"text": value, "size": 10, "bold": value in {"PASS", "FAIL", "FAIL*"}, "color": color, "align": PP_ALIGN.CENTER if x != 2.0 else PP_ALIGN.LEFT}], margin=0.03)
        y += row_h + 0.02
    add_box_with_text(slide, 0.85, 6.55, 11.8, 0.36, fill=WARNING_BG, line=AMBER, lines=[{"text": "* bd67944a fixed the bug but produced unparseable evaluator JSON; the prompt was hardened afterwards.", "size": 11, "color": NAVY}], margin=0.03)



def slide_a2(slide, number: int) -> None:
    set_background(slide, LIGHT_BG)
    add_title(slide, number, "Appendix — Evaluator Reliability Metrics")
    add_metric_card(slide, 0.95, 1.75, 2.5, 1.0, "Rubber-stamp incidents", "0 explicit", fill=SUCCESS_BG, accent=TEAL)
    add_metric_card(slide, 3.75, 1.75, 2.5, 1.0, "Refusal-to-be-negative", "1 incident", fill=DANGER_BG, accent=CORAL)
    add_metric_card(slide, 6.55, 1.75, 2.6, 1.0, "False positive rate", "Deferred", fill=WHITE, accent=GRAY)
    add_metric_card(slide, 9.45, 1.75, 2.3, 1.0, "Cross-model agreement", "Not measured", fill=WHITE, accent=GRAY)
    add_box_with_text(slide, 0.95, 3.3, 5.45, 2.1, fill=WHITE, line=CORAL, lines=[{"text": "Closest miss", "size": 19, "bold": True, "color": CORAL}, {"text": "TS benchmark: evaluator gave functionality = 8 even though only 11/22 tests passed. That is exactly the people-pleasing bias the article warned about.", "size": 15, "color": NAVY}], margin=0.08)
    add_box_with_text(slide, 6.85, 3.3, 5.1, 2.1, fill=WHITE, line=TEAL, lines=[{"text": "Best evidence", "size": 19, "bold": True, "color": TEAL_DARK}, {"text": "Python bugfix: functionality went 1 → 10 after evaluator feedback. The critic was harsh first, then precise enough to drive the fix.", "size": 15, "color": NAVY}], margin=0.08)
    add_box_with_text(slide, 1.2, 6.0, 10.8, 0.52, fill=NAVY, line=NAVY, lines=[{"text": "Calibration solved format problems fast. Severity scoring remains the harder problem.", "size": 14, "color": WHITE, "align": PP_ALIGN.CENTER}], margin=0.04)



def slide_a3(slide, number: int) -> None:
    set_background(slide, LIGHT_BG)
    add_title(slide, number, "Appendix — Difficulty Classification")
    x, y, w, h = 1.2, 1.7, 7.0, 4.2
    add_line(slide, x, y + h, x + w, y + h, color=NAVY, width=1.6)
    add_line(slide, x, y, x, y + h, color=NAVY, width=1.6)
    add_line(slide, x + w / 2, y, x + w / 2, y + h, color=GRAY, width=1.0)
    add_line(slide, x, y + h / 2, x + w, y + h / 2, color=GRAY, width=1.0)
    add_textbox(slide, x + 1.7, y + h + 0.08, 3.5, 0.22, [{"text": "Higher solo capability →", "size": 12, "color": NAVY, "align": PP_ALIGN.CENTER}], margin=0)
    add_textbox(slide, x - 0.95, y + 1.4, 0.8, 0.8, [{"text": "Trio uplift", "size": 12, "color": NAVY, "align": PP_ALIGN.CENTER}], margin=0)
    points = [
        (6.2, 2.6, "Python bugfix", GRAY, "marginal"),
        (6.1, 2.2, "TS feature", GRAY, "marginal"),
        (4.3, 1.4, "Go race", CORAL, "too_hard"),
        (6.1, 2.0, "Python tags", GRAY, "marginal"),
        (4.9, 4.2, "Fullstack", TEAL, "in_zone"),
    ]
    for solo, trio_y, label, color, classification in points:
        plot_x = x + (solo / 9.5) * w
        plot_y = y + h - (trio_y / 4.5) * h
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(plot_x - 0.11), Inches(plot_y - 0.11), Inches(0.22), Inches(0.22))
        dot.fill.solid(); dot.fill.fore_color.rgb = color; dot.line.color.rgb = WHITE
        add_textbox(slide, plot_x + 0.1, plot_y - 0.08, 1.0, 0.2, [{"text": label, "size": 10, "color": NAVY}], margin=0)
        add_pill(slide, plot_x - 0.2, plot_y + 0.12, 0.75, 0.18, classification, fill=color if color != GRAY else NAVY, color=WHITE, size=8)
    add_box_with_text(slide, 8.75, 1.8, 3.65, 3.6, fill=WHITE, line=NEUTRAL_BG, lines=[
        {"text": "Legend", "size": 18, "bold": True, "color": NAVY},
        {"text": "too_easy: both ≥ 9", "size": 14, "color": NAVY},
        {"text": "too_hard: both fail", "size": 14, "color": NAVY},
        {"text": "in_zone: trio wins by ≥ 1.5", "size": 14, "color": NAVY},
        {"text": "marginal: small uplift or tie", "size": 14, "color": NAVY},
        {"text": "Observation: only the fullstack benchmark landed squarely in the harness sweet spot.", "size": 14, "bold": True, "color": AMBER, "space_after": 0},
    ], margin=0.08)



def slide_a4(slide, number: int) -> None:
    set_background(slide, LIGHT_BG)
    add_title(slide, number, "Appendix — Grading Criteria Details")
    add_box_with_text(slide, 0.85, 1.45, 5.8, 4.95, fill=WHITE, line=TEAL, lines=[{"text": "backend.yaml", "size": 22, "bold": True, "color": TEAL_DARK, "align": PP_ALIGN.CENTER}], margin=0.08)
    y = 2.0
    for name, weight, threshold, desc in BACKEND_CRITERIA:
        add_box_with_text(slide, 1.05, y, 5.4, 0.78, fill=SUCCESS_BG if weight == "HIGH" else LIGHT_PANEL, line=TEAL if weight == "HIGH" else GRAY, lines=[{"text": f"{name} · {weight} · threshold {threshold}", "size": 14, "bold": True, "color": NAVY}, {"text": desc, "size": 11, "color": NAVY}], margin=0.05)
        y += 0.9
    add_box_with_text(slide, 6.9, 1.45, 5.6, 4.95, fill=WHITE, line=PURPLE, lines=[{"text": "fullstack.yaml", "size": 22, "bold": True, "color": PURPLE, "align": PP_ALIGN.CENTER}], margin=0.08)
    y = 2.0
    for name, weight, threshold, desc in FULLSTACK_CRITERIA:
        add_box_with_text(slide, 7.1, y, 5.2, 0.78, fill=LAVENDER if name == "Visual Design" else SUCCESS_BG if weight == "HIGH" else LIGHT_PANEL, line=PURPLE if name == "Visual Design" else TEAL if weight == "HIGH" else GRAY, lines=[{"text": f"{name} · {weight} · threshold {threshold}", "size": 14, "bold": True, "color": NAVY}, {"text": desc, "size": 11, "color": NAVY}], margin=0.05)
        y += 0.9



def slide_a5(slide, number: int) -> None:
    set_background(slide, NAVY)
    add_title(slide, number, "Appendix — Detailed System Diagram", dark=True)
    add_box_with_text(slide, 0.85, 1.45, 1.8, 0.8, fill=AMBER, line=AMBER, lines=[{"text": "Human prompt", "size": 18, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER}], margin=0.05)
    add_box_with_text(slide, 3.0, 1.25, 2.3, 1.0, fill=BLUE_GRAY, line=BLUE_GRAY, lines=[{"text": "Planner", "size": 20, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}, {"text": "spec.md + design direction", "size": 13, "color": WHITE, "align": PP_ALIGN.CENTER}], margin=0.05)
    add_box_with_text(slide, 3.0, 3.0, 2.3, 1.0, fill=TEAL, line=TEAL, lines=[{"text": "Generator", "size": 20, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}, {"text": "builds against agreed sprint contract", "size": 13, "color": WHITE, "align": PP_ALIGN.CENTER}], margin=0.05)
    add_box_with_text(slide, 6.0, 3.0, 2.3, 1.0, fill=CORAL, line=CORAL, lines=[{"text": "Evaluator", "size": 20, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}, {"text": "Playwright QA + scoring + bug reports", "size": 13, "color": WHITE, "align": PP_ALIGN.CENTER}], margin=0.05)
    add_box_with_text(slide, 9.0, 1.25, 3.1, 1.0, fill=GRAY_DARK, line=GRAY_DARK, lines=[{"text": "Telemetry collector", "size": 20, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}, {"text": "timing · scores · bugs · trends · JSON", "size": 13, "color": WHITE, "align": PP_ALIGN.CENTER}], margin=0.05)
    add_box_with_text(slide, 9.0, 3.0, 3.1, 1.0, fill=NAVY_DARK, line=TEAL, lines=[{"text": "Replay / reporting", "size": 20, "bold": True, "color": TEAL, "align": PP_ALIGN.CENTER}, {"text": "re-score saved artifacts without rerunning the generator", "size": 13, "color": WHITE, "align": PP_ALIGN.CENTER}], margin=0.05)
    add_box_with_text(slide, 0.95, 5.05, 3.15, 1.1, fill=NAVY_DARK, line=AMBER, lines=[{"text": "Isolation manager", "size": 19, "bold": True, "color": AMBER, "align": PP_ALIGN.CENTER}, {"text": "Generator cannot see _eval/ hidden acceptance tests", "size": 13, "color": WHITE, "align": PP_ALIGN.CENTER}], margin=0.05)
    add_box_with_text(slide, 4.55, 5.05, 3.15, 1.1, fill=NAVY_DARK, line=TEAL, lines=[{"text": "Communication protocol", "size": 19, "bold": True, "color": TEAL, "align": PP_ALIGN.CENTER}, {"text": "Everything is a file handoff — no shared chat memory", "size": 13, "color": WHITE, "align": PP_ALIGN.CENTER}], margin=0.05)
    add_box_with_text(slide, 8.1, 5.05, 4.0, 1.1, fill=NAVY_DARK, line=CORAL, lines=[{"text": "Context strategy", "size": 19, "bold": True, "color": CORAL, "align": PP_ALIGN.CENTER}, {"text": "Reset between phases when models exhibit context anxiety; compact when they do not", "size": 13, "color": WHITE, "align": PP_ALIGN.CENTER}], margin=0.05)
    add_line(slide, 2.65, 1.85, 2.95, 1.75, color=WHITE, width=2.0)
    add_line(slide, 4.15, 2.25, 4.15, 2.95, color=WHITE, width=2.0)
    add_line(slide, 5.3, 3.5, 5.95, 3.5, color=AMBER, width=2.2)
    add_line(slide, 8.35, 3.5, 9.0, 3.5, color=WHITE, width=2.0)
    add_line(slide, 5.3, 1.75, 8.95, 1.75, color=WHITE, width=2.0)


SLIDE_BUILDERS = [
    slide_1,
    slide_2,
    slide_3,
    slide_4,
    slide_5,
    slide_6,
    slide_7,
    slide_8,
    slide_9,
    slide_10,
    slide_11,
    slide_12,
    slide_13,
    slide_14,
    slide_15,
    slide_16,
    slide_17,
    slide_18,
    slide_19,
    slide_20,
    slide_21,
    slide_22,
    slide_a1,
    slide_a2,
    slide_a3,
    slide_a4,
    slide_a5,
]



def build_presentation(output_path: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = SLIDE_WIDTH
    presentation.slide_height = SLIDE_HEIGHT
    presentation.core_properties.title = TITLE
    presentation.core_properties.subject = "Engineering leadership talk deck"
    presentation.core_properties.author = "GitHub Copilot CLI"
    presentation.core_properties.keywords = "AI agents, adversarial architecture, harnessa"

    for index, builder in enumerate(SLIDE_BUILDERS, start=1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        builder(slide, index)

    presentation.save(output_path)



def validate_presentation(output_path: Path) -> None:
    reopened = Presentation(output_path)
    if len(reopened.slides) != TOTAL_SLIDES:
        raise ValueError(f"expected {TOTAL_SLIDES} slides, found {len(reopened.slides)}")
    with zipfile.ZipFile(output_path) as archive:
        slide_xml = [name for name in archive.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")]
    if len(slide_xml) != TOTAL_SLIDES:
        raise ValueError(f"pptx contains {len(slide_xml)} slide XML files, expected {TOTAL_SLIDES}")



def generate_quicklook_preview(output_path: Path, preview_dir: Path) -> Path | None:
    qlmanage = shutil.which("qlmanage")
    if qlmanage is None:
        return None

    bundle_dir = preview_dir / f"{output_path.name}.qlpreview"
    if bundle_dir.exists():
        shutil.rmtree(bundle_dir)

    result = subprocess.run(
        [qlmanage, "-p", "-o", str(preview_dir), str(output_path)],
        capture_output=True,
        text=True,
        check=False,
    )
    if result.returncode != 0:
        return None

    preview_html = bundle_dir / "Preview.html"
    return preview_html if preview_html.exists() else None


def build_fallback_preview(output_path: Path, fallback_path: Path) -> None:
    presentation = Presentation(output_path)
    sections: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts = [
            shape.text.strip().replace("\n", " ")
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        title = texts[0] if texts else f"Slide {index}"
        bullets = texts[1:6]
        bullet_html = "".join(f"<li>{escape(text)}</li>" for text in bullets)
        theme = "dark" if index in FALLBACK_DARK_SLIDES else "light"
        sections.append(
            f"""
            <section class="slide {theme}">
              <div class="slide-number">{index:02d}</div>
              <h2>{escape(title)}</h2>
              <ul>{bullet_html}</ul>
            </section>
            """
        )

    fallback_path.write_text(
        f"""<!doctype html>
<html lang="en">
  <head>
    <meta charset="utf-8">
    <meta name="viewport" content="width=device-width, initial-scale=1">
    <title>{escape(TITLE)} — fallback preview</title>
    <style>
      :root {{
        --navy: #1b2a4a;
        --panel: #ffffff;
        --ink: #1b2a4a;
        --teal: #2ec4b6;
        --shadow: 0 18px 40px rgba(0, 0, 0, 0.12);
      }}
      * {{ box-sizing: border-box; }}
      body {{
        margin: 0;
        background: linear-gradient(180deg, #eef2f7 0%, #dbe4ef 100%);
        color: var(--ink);
        font-family: Calibri, "Segoe UI", Arial, sans-serif;
      }}
      main {{
        max-width: 1480px;
        margin: 0 auto;
        padding: 24px;
      }}
      .hero {{
        margin-bottom: 24px;
        padding: 20px 24px;
        border-radius: 20px;
        background: var(--navy);
        color: white;
        box-shadow: var(--shadow);
      }}
      .hero p {{ margin: 8px 0 0; color: #d8fffb; }}
      .grid {{
        display: grid;
        grid-template-columns: repeat(auto-fit, minmax(440px, 1fr));
        gap: 24px;
      }}
      .slide {{
        position: relative;
        min-height: 360px;
        padding: 32px 36px;
        border-radius: 18px;
        box-shadow: var(--shadow);
      }}
      .slide.light {{
        background: var(--panel);
        color: var(--ink);
      }}
      .slide.dark {{
        background: var(--navy);
        color: white;
      }}
      .slide.dark h2 {{
        color: white;
      }}
      .slide.dark li {{
        color: #ebffff;
      }}
      .slide-number {{
        position: absolute;
        top: 18px;
        right: 22px;
        font-weight: 700;
        letter-spacing: 0.12em;
        color: var(--teal);
      }}
      h2 {{
        margin: 0 0 16px;
        font-size: 30px;
        line-height: 1.15;
      }}
      ul {{
        margin: 0;
        padding-left: 22px;
        line-height: 1.45;
        font-size: 19px;
      }}
      li + li {{ margin-top: 10px; }}
    </style>
  </head>
  <body>
    <main>
      <section class="hero">
        <h1>{escape(TITLE)} — fallback browser preview</h1>
        <p>Generated from the PPTX when Quick Look HTML export is unavailable. It preserves slide order, titles, and key text for browser inspection.</p>
      </section>
      <section class="grid">
        {''.join(sections)}
      </section>
    </main>
  </body>
</html>
""",
        encoding="utf-8",
    )


def build_preview_index(
    preview_dir: Path,
    output_path: Path,
    plan_copy: Path,
    quicklook_preview: Path | None,
) -> Path:
    preview_dir.mkdir(exist_ok=True)
    fallback_path = preview_dir / "fallback-preview.html"
    build_fallback_preview(output_path, fallback_path)

    iframe_src = (
        escape(quicklook_preview.relative_to(preview_dir).as_posix())
        if quicklook_preview is not None
        else escape(fallback_path.name)
    )
    quicklook_link = (
        f'<a href="{escape(quicklook_preview.relative_to(preview_dir).as_posix())}">Quick Look HTML preview</a>'
        if quicklook_preview is not None
        else '<span>Quick Look HTML preview unavailable on this machine</span>'
    )
    pdf_path = output_path.with_suffix(".pdf")
    pdf_link = (
        f'<a class="secondary" href="../{escape(pdf_path.name)}">PDF export</a>'
        if pdf_path.exists()
        else ""
    )
    output_rel = escape(f"../{output_path.name}")
    plan_rel = escape(f"../{plan_copy.name}")
    fallback_rel = escape(fallback_path.name)

    index_path = preview_dir / "index.html"
    index_path.write_text(
        f"""<!doctype html>
<html lang="en">
  <head>
    <meta charset="utf-8">
    <meta name="viewport" content="width=device-width, initial-scale=1">
    <title>{escape(TITLE)} — preview</title>
    <style>
      body {{
        margin: 0;
        background: #0f1724;
        color: #f8f8f8;
        font-family: Calibri, "Segoe UI", Arial, sans-serif;
      }}
      .shell {{
        display: grid;
        grid-template-rows: auto 1fr;
        min-height: 100vh;
      }}
      header {{
        padding: 20px 24px 14px;
        border-bottom: 1px solid rgba(255, 255, 255, 0.15);
        background: rgba(19, 32, 58, 0.95);
      }}
      h1 {{
        margin: 0;
        font-size: 28px;
      }}
      p {{
        margin: 8px 0 14px;
        color: #c6fff9;
      }}
      nav {{
        display: flex;
        flex-wrap: wrap;
        gap: 12px;
      }}
      nav a, nav span {{
        color: #0f1724;
        background: #f2cc8f;
        text-decoration: none;
        padding: 10px 14px;
        border-radius: 999px;
        font-weight: 700;
      }}
      nav a.secondary {{
        background: #2ec4b6;
      }}
      .frame {{
        padding: 18px;
      }}
      iframe {{
        width: 100%;
        min-height: calc(100vh - 150px);
        border: 0;
        border-radius: 18px;
        background: white;
      }}
    </style>
  </head>
  <body>
    <div class="shell">
      <header>
        <h1>{escape(TITLE)} — browser preview</h1>
        <p>Open this page in a browser-style tool to inspect rendered slides. The frame below prefers the Quick Look HTML export and falls back to a generated HTML summary.</p>
        <nav>
          <a href="{output_rel}">Deck (.pptx)</a>
          {pdf_link}
          <a class="secondary" href="{plan_rel}">Plan copy</a>
          {quicklook_link}
          <a class="secondary" href="{fallback_rel}">Fallback HTML preview</a>
        </nav>
      </header>
      <div class="frame">
        <iframe src="{iframe_src}" title="{escape(TITLE)} preview"></iframe>
      </div>
    </div>
  </body>
</html>
""",
        encoding="utf-8",
    )
    return index_path


def validate_preview(preview_dir: Path, preview_index: Path, quicklook_preview: Path | None) -> None:
    if not preview_index.exists():
        raise ValueError("preview index was not generated")
    fallback_preview = preview_dir / "fallback-preview.html"
    if not fallback_preview.exists():
        raise ValueError("fallback preview was not generated")
    if quicklook_preview is not None:
        if not quicklook_preview.exists():
            raise ValueError("quicklook preview html is missing")


def main() -> None:
    repo_root = Path(__file__).resolve().parents[1]
    presentation_dir = repo_root / "presentation"
    presentation_dir.mkdir(exist_ok=True)
    output_path = presentation_dir / "the-adversarial-architecture.pptx"
    pdf_path = output_path.with_suffix(".pdf")
    plan_source = repo_root / "docs" / "PRESENTATION_PLAN.md"
    plan_copy = presentation_dir / "PRESENTATION_PLAN.md"
    preview_dir = presentation_dir / "preview"

    build_presentation(output_path)
    validate_presentation(output_path)
    shutil.copy2(plan_source, plan_copy)
    preview_dir.mkdir(exist_ok=True)
    quicklook_preview = generate_quicklook_preview(output_path, preview_dir)
    preview_index = build_preview_index(preview_dir, output_path, plan_copy, quicklook_preview)
    validate_preview(preview_dir, preview_index, quicklook_preview)

    print(f"Built {output_path.relative_to(repo_root)} ({TOTAL_SLIDES} slides).")
    print(f"Copied {plan_copy.relative_to(repo_root)}.")
    if pdf_path.exists():
        print(f"PDF export: {pdf_path.relative_to(repo_root)}.")
    print(f"Preview index: {preview_index.relative_to(repo_root)}.")


if __name__ == "__main__":
    main()
